"""Parse various file formats to plain text."""

import io
from pathlib import Path


def parse_file(file_name: str, file_bytes: bytes) -> tuple[str, str]:
    """
    Parse a file and return (plain_text_content, file_type).
    Supported: PDF, DOCX, PPTX, TXT, MD, CSV, XLSX
    """
    name_lower = file_name.lower()
    ext = Path(name_lower).suffix

    # Determine type
    if ext == ".pdf":
        return _parse_pdf(file_bytes), "pdf"
    elif ext == ".docx":
        return _parse_docx(file_bytes), "docx"
    elif ext == ".pptx":
        return _parse_pptx(file_bytes), "pptx"
    elif ext in (".txt", ".md", ".csv", ".log", ".json", ".xml", ".yaml", ".yml", ".py", ".js", ".ts", ".html", ".css"):
        return _parse_text(file_bytes), ext.lstrip(".")
    elif ext in (".xlsx", ".xls"):
        return _parse_xlsx(file_bytes), "xlsx"
    else:
        # Try text fallback
        try:
            return _parse_text(file_bytes), "txt"
        except Exception:
            raise ValueError(f"Unsupported file format: {ext}")


def _parse_pdf(data: bytes) -> str:
    """Extract text from PDF."""
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    parts: list[str] = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            parts.append(text)
    return "\n\n".join(parts) if parts else "(PDF contained no extractable text)"


def _parse_docx(data: bytes) -> str:
    """Extract text from DOCX."""
    from docx import Document
    doc = Document(io.BytesIO(data))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    # Also extract tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)
    return "\n".join(parts) if parts else "(DOCX contained no extractable text)"


def _parse_pptx(data: bytes) -> str:
    """Extract text from PPTX slides."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        slide_texts.append(paragraph.text.strip())
        if slide_texts:
            parts.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_texts))
    return "\n\n".join(parts) if parts else "(PPTX contained no extractable text)"


def _parse_xlsx(data: bytes) -> str:
    """Extract text from XLSX spreadsheets."""
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = [f"--- Sheet: {sheet_name} ---"]
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(c) if c is not None else "" for c in row)
            if row_text.strip():
                rows.append(row_text)
        parts.append("\n".join(rows))
    wb.close()
    return "\n\n".join(parts) if parts else "(XLSX contained no data)"


def _parse_text(data: bytes) -> str:
    """Read text file with encoding detection."""
    # Try UTF-8 first, then common fallbacks
    for enc in ("utf-8", "gbk", "gb2312", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")
